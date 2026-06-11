#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
城发投经营财务月报 - PowerPoint简报模板生成器
自动生成包含11页幻灯片的PPT模板
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_ppt_template():
    """创建PPT简报模板"""
    
    # 创建演示文稿（16:9宽屏）
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 定义配色方案
    COLOR_PRIMARY = RGBColor(0, 82, 204)      # 蓝色 #0052CC
    COLOR_SECONDARY = RGBColor(255, 107, 53)   # 橙色 #FF6B35
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_BLACK = RGBColor(0, 0, 0)
    COLOR_GRAY = RGBColor(100, 100, 100)
    
    # ========== 第1页：封面 ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 背景色（蓝色渐变效果用纯色代替）
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY
    
    # 标题
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "城发投经营财务月报"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.333), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "2026年X月"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = COLOR_WHITE
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # 编制单位
    unit_box = slide1.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(0.5))
    unit_frame = unit_box.text_frame
    unit_frame.text = "编制单位：综合管理部"
    unit_para = unit_frame.paragraphs[0]
    unit_para.font.size = Pt(18)
    unit_para.font.color.rgb = COLOR_WHITE
    unit_para.alignment = PP_ALIGN.CENTER
    
    # 日期
    date_box = slide1.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11.333), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = "编制日期：2026年X月X日"
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(18)
    date_para.font.color.rgb = COLOR_WHITE
    date_para.alignment = PP_ALIGN.CENTER
    
    # 本期亮点（预留位置）
    highlight_box = slide1.shapes.add_textbox(Inches(1), Inches(7.0), Inches(11.333), Inches(0.5))
    highlight_frame = highlight_box.text_frame
    highlight_frame.text = "本期亮点：（每月填写）"
    highlight_para = highlight_frame.paragraphs[0]
    highlight_para.font.size = Pt(16)
    highlight_para.font.color.rgb = COLOR_WHITE
    highlight_para.alignment = PP_ALIGN.CENTER
    
    # ========== 第2页：目录 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "目录"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRIMARY
    
    # 目录内容
    content_box = slide2.shapes.add_textbox(Inches(2), Inches(1.8), Inches(9.333), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.text = """一、本月经营概览

二、PPP项目板块

三、供热项目板块

四、平台公司板块

五、本部板块

六、风险预警与建议"""
    for para in content_frame.paragraphs:
        para.font.size = Pt(24)
        para.font.color.rgb = COLOR_BLACK
        para.space_after = Pt(15)
    
    # ========== 第3页：本月经营概览 ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "一、本月经营概览"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRIMARY
    
    # 核心指标卡片（左上）
    card1 = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(2.5))
    tf = card1.text_frame
    tf.text = """核心指标

• 营业收入：18.82亿元 ▲+5.2%
• 利润总额：1.65亿元 ▼-10.3%
• 经营性现金流：3.27亿元 ▲+15.6%
• 净资产收益率：3.53% ▼-0.8pp"""
    for para in tf.paragraphs:
        if para == tf.paragraphs[0]:
            para.font.size = Pt(20)
            para.font.bold = True
            para.font.color.rgb = COLOR_PRIMARY
        else:
            para.font.size = Pt(16)
            para.font.color.rgb = COLOR_BLACK
            para.space_after = Pt(8)
    
    # 本月亮点（中右）
    highlight = slide3.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(6), Inches(1.2))
    tf = highlight.text_frame
    tf.text = """本月亮点：

1. PPP项目回款创历史新高
2. 供热项目成本控制显著
3. 平台公司管理提升明显"""
    for para in tf.paragraphs:
        if para == tf.paragraphs[0]:
            para.font.size = Pt(18)
            para.font.bold = True
            para.font.color.rgb = COLOR_SECONDARY
        else:
            para.font.size = Pt(14)
            para.font.color.rgb = COLOR_BLACK
            para.space_after = Pt(6)
    
    # 本月关注（底部）
    concern = slide3.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(2))
    tf = concern.text_frame
    tf.text = """本月关注：

1. PPP项目应收账款余额较高
2. 供热项目用户满意度下降
3. 本部管理费用控制需加强"""
    for para in tf.paragraphs:
        if para == tf.paragraphs[0]:
            para.font.size = Pt(18)
            para.font.bold = True
            para.font.color.rgb = COLOR_SECONDARY
        else:
            para.font.size = Pt(14)
            para.font.color.rgb = COLOR_BLACK
            para.space_after = Pt(6)
    
    # 图表占位符说明（右上和左下）
    note1 = slide3.shapes.add_textbox(Inches(6.5), Inches(2.8), Inches(6), Inches(0.5))
    tf = note1.text_frame
    tf.text = "[近12个月趋势图 - 折线图占位符]"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = COLOR_GRAY
    tf.paragraphs[0].font.italic = True
    
    note2 = slide3.shapes.add_textbox(Inches(0.5), Inches(4), Inches(5.5), Inches(0.4))
    tf = note2.text_frame
    tf.text = "[板块利润贡献饼图占位符]"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = COLOR_GRAY
    tf.paragraphs[0].font.italic = True
    
    # ========== 第4页：PPP项目板块-核心指标 ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "二、PPP项目板块（1/2）"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRIMARY
    
    # 核心指标表格占位
    table_note = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(0.8))
    tf = table_note.text_frame
    tf.text = "[核心指标表格占位符 - 包含营业收入、利润总额、回款金额、应收账款余额等指标]"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLOR_GRAY
    tf.paragraphs[0].font.italic = True
    
    # 对比分析占位
    compare_note = slide4.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(0.8))
    tf = compare_note.text_frame
    tf.text = "[同比/环比/年度目标对比分析占位符]"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLOR_GRAY
    tf.paragraphs[0].font.italic = True
    
    # 项目状态占位
    status_note = slide4.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12), Inches(0.8))
    tf = status_note.text_frame
    tf.text = "[项目状态分布饼图占位符 - 正常/预警/问题项目]"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLOR_GRAY
    tf.paragraphs[0].font.italic = True
    
    # ========== 第5页：PPP项目板块-效益排行 ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "二、PPP项目板块（2/2）"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRIMARY
    
    # 效益排行占位
    ranking_note = slide5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(2))
    tf = ranking_note.text_frame
    tf.text = """[项目效益排行占位符]

利润贡献前5名项目：
1. 许昌项目：0.25亿元
2. 南阳项目：0.18亿元
3. 安阳项目：0.15亿元
4. 鹤壁项目：0.12亿元
5. 濮阳项目：0.10亿元

利润贡献后5名项目：
1. 信阳项目：-0.05亿元
2. 漯河项目：-0.02亿元
..."""
    for para in tf.paragraphs:
        para.font.size = Pt(14)
        para.font.color.rgb = COLOR_BLACK
        para.space_after = Pt(6)
    
    # 回款情况占位
    payment_note = slide5.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12), Inches(1.5))
    tf = payment_note.text_frame
    tf.text = "[回款情况分析占位符 - 包含回款金额、回款率、账龄分析等]"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLOR_GRAY
    tf.paragraphs[0].font.italic = True
    
    # 问题跟踪占位
    issue_note = slide5.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(1.5))
    tf = issue_note.text_frame
    tf.text = "[问题项目跟踪占位符 - 包含问题原因、改进措施、责任部门等]"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLOR_GRAY
    tf.paragraphs[0].font.italic = True
    
    # ========== 第6页：供热项目板块-核心指标 ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide6.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "三、供热项目板块（1/2）"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRIMARY
    
    # 核心指标占位
    note = slide6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(0.8))
    tf = note.text_frame
    tf.text = "[核心指标表格占位符 - 包含供热面积、营业收入、利润总额、单位面积利润等]"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLOR_GRAY
    tf.paragraphs[0].font.italic = True
    
    # 区域分布占位
    map_note = slide6.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(0.8))
    tf = map_note.text_frame
    tf.text = "[区域分布地图热力图占位符 - 显示各区域供热面积和利润贡献]"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLOR_GRAY
    tf.paragraphs[0].font.italic = True
    
    # ========== 第7页：供热项目板块-成本分析 ==========
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide7.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "三、供热项目板块（2/2）"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRIMARY
    
    # 成本分析占位
    cost_note = slide7.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(0.8))
    tf = cost_note.text_frame
    tf.text = "[成本构成饼图占位符 - 燃料成本60%、人工成本25%、维护成本10%、其他5%]"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLOR_GRAY
    tf.paragraphs[0].font.italic = True
    
    # 效率分析占位
    eff_note = slide7.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(0.8))
    tf = eff_note.text_frame
    tf.text = "[效率分析占位符 - 单位面积成本、能耗比、人均供热面积等]"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLOR_GRAY
    tf.paragraphs[0].font.italic = True
    
    # 用户满意度占位
    satis_note = slide7.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12), Inches(0.8))
    tf = satis_note.text_frame
    tf.text = "[用户满意度分析占位符 - 满意度评分、投诉率、响应时间等]"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLOR_GRAY
    tf.paragraphs[0].font.italic = True
    
    # ========== 第8页：平台公司板块-核心指标 ==========
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide8.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "四、平台公司板块（1/2）"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRIMARY
    
    # 核心指标占位
    note = slide8.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(0.8))
    tf = note.text_frame
    tf.text = "[核心指标表格占位符 - 包含公司数量、营业收入、利润总额、净资产等]"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLOR_GRAY
    tf.paragraphs[0].font.italic = True
    
    # 公司状态占位
    status_note = slide8.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(0.8))
    tf = status_note.text_frame
    tf.text = "[公司状态分布占位符 - 盈利/盈亏平衡/亏损公司数量及占比]"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLOR_GRAY
    tf.paragraphs[0].font.italic = True
    
    # ========== 第9页：平台公司板块-效益排行 ==========
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide9.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "四、平台公司板块（2/2）"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRIMARY
    
    # 效益排行占位
    ranking_note = slide9.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(2))
    tf = ranking_note.text_frame
    tf.text = """[公司效益排行占位符]

利润贡献前3名公司：
1. 新郑公司：0.10亿元
2. 南阳公司：0.08亿元
3. 漯河公司：0.05亿元

利润贡献后3名公司：
1. 鄢陵公司：-0.01亿元
2. 启迪零碳：0.01亿元
3. 港区运营：0.02亿元"""
    for para in tf.paragraphs:
        para.font.size = Pt(14)
        para.font.color.rgb = COLOR_BLACK
        para.space_after = Pt(6)
    
    # 资产质量占位
    asset_note = slide9.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12), Inches(1.5))
    tf = asset_note.text_frame
    tf.text = "[资产质量分析占位符 - 资产负债率、净资产收益率、资产周转率等]"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLOR_GRAY
    tf.paragraphs[0].font.italic = True
    
    # 问题跟踪占位
    issue_note = slide9.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(1.5))
    tf = issue_note.text_frame
    tf.text = "[问题公司跟踪占位符 - 包含亏损原因、改进措施、责任部门等]"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLOR_GRAY
    tf.paragraphs[0].font.italic = True
    
    # ========== 第10页：本部板块 ==========
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide10.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "五、本部板块"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRIMARY
    
    # 核心指标占位
    note = slide10.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(0.8))
    tf = note.text_frame
    tf.text = "[核心指标表格占位符 - 包含管理费用、财务费用、人员数量、人均费用等]"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLOR_GRAY
    tf.paragraphs[0].font.italic = True
    
    # 费用控制占位
    cost_note = slide10.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(0.8))
    tf = cost_note.text_frame
    tf.text = "[费用构成饼图占位符 - 职工薪酬60%、办公费15%、业务招待费10%、其他15%]"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLOR_GRAY
    tf.paragraphs[0].font.italic = True
    
    # 与预算对比占位
    budget_note = slide10.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12), Inches(1.5))
    tf = budget_note.text_frame
    tf.text = "[与预算对比占位符 - 职工薪酬节约5.2%、办公费超预算3.1%、业务招待费节约12.5%等]"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLOR_GRAY
    tf.paragraphs[0].font.italic = True
    
    # 效率提升占位
    eff_note = slide10.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(12), Inches(1.5))
    tf = eff_note.text_frame
    tf.text = "[效率提升分析占位符 - 人均处理业务量、流程优化效果、数字化程度等]"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLOR_GRAY
    tf.paragraphs[0].font.italic = True
    
    # ========== 第11页：风险预警与建议 ==========
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide11.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "六、风险预警与建议"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRIMARY
    
    # 风险预警列表
    risk_box = slide11.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(2))
    tf = risk_box.text_frame
    tf.text = """风险预警：

🔴 高风险：PPP项目应收账款余额较高（58.84亿元，回收周期长）
🟡 中风险：供热项目用户满意度下降（评分下降0.5分）
🟢 低风险：本部管理费用控制需加强（部分费用超预算）"""
    for para in tf.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = COLOR_BLACK
        para.space_after = Pt(10)
    
    # 建议措施
    suggestion_box = slide11.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12), Inches(1.5))
    tf = suggestion_box.text_frame
    tf.text = """建议措施：

1. 加大PPP项目回款力度，建立双周例会制度
2. 改进供热项目服务质量，加强用户沟通
3. 加强本部费用控制，严格执行预算
4. 推进轻资产项目落地，提升资产收益率"""
    for para in tf.paragraphs:
        if para == tf.paragraphs[0]:
            para.font.size = Pt(18)
            para.font.bold = True
            para.font.color.rgb = COLOR_SECONDARY
        else:
            para.font.size = Pt(14)
            para.font.color.rgb = COLOR_BLACK
            para.space_after = Pt(8)
    
    # 下月关注重点
    next_month_box = slide11.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(1.5))
    tf = next_month_box.text_frame
    tf.text = """下月关注重点：

• PPP项目：许昌、长垣、濮阳等5大攻坚目标回款进度
• 供热项目：2025-2026采暖季收尾工作及成本结算
• 平台公司：新郑、南阳等公司利润完成情况
• 本部：管理费用预算执行及全年压降目标"""
    for para in tf.paragraphs:
        para.font.size = Pt(14)
        para.font.color.rgb = COLOR_BLACK
        para.space_after = Pt(6)
    
    # 保存文件
    output_file = "城发投经营财务月报模板.pptx"
    prs.save(output_file)
    print(f"✅ PPT模板已生成：{output_file}")
    print(f"📊 共生成 {len(prs.slides)} 页幻灯片")
    print("\n📝 每页内容说明：")
    print("   第1页：封面")
    print("   第2页：目录")
    print("   第3页：本月经营概览")
    print("   第4-5页：PPP项目板块（2页）")
    print("   第6-7页：供热项目板块（2页）")
    print("   第8-9页：平台公司板块（2页）")
    print("   第10页：本部板块")
    print("   第11页：风险预警与建议")
    print("\n💡 使用说明：")
    print("   1. 打开生成的.pptx文件")
    print("   2. 替换占位符内容为实际数据和图表")
    print("   3. 根据实际需求调整配色和布局")
    print("   4. 另存为.potx文件作为模板使用")

if __name__ == "__main__":
    create_ppt_template()
